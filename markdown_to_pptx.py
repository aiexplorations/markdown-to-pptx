#!/usr/bin/env python3
"""
Markdown to PPTX Converter - POC

A simple proof of concept script to convert Markdown files to PowerPoint presentations.
Supports basic markdown elements like headers, bullet points, and text.
"""

import argparse
import re
from pathlib import Path
from typing import List, Dict, Any, Tuple
import urllib.request
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    exit(1)


class MarkdownToPPTX:
    """Convert Markdown content to PowerPoint presentation."""
    
    def __init__(self) -> None:
        """Initialize the converter."""
        self.prs = Presentation()
        self.current_slide = None
        
    def parse_text_formatting(self, text: str) -> List[Tuple[str, bool, bool]]:
        """Parse text and return segments with formatting info (text, is_bold, is_italic)."""
        segments = []
        i = 0
        current_text = ""
        
        while i < len(text):
            # Check for bold (**text** or __text__)
            if (i < len(text) - 1 and text[i:i+2] == '**') or (i < len(text) - 1 and text[i:i+2] == '__'):
                if current_text:
                    segments.append((current_text, False, False))
                    current_text = ""
                
                delimiter = text[i:i+2]
                i += 2
                bold_text = ""
                
                # Find closing delimiter
                while i < len(text) - 1:
                    if text[i:i+2] == delimiter:
                        segments.append((bold_text, True, False))
                        i += 2
                        break
                    bold_text += text[i]
                    i += 1
                else:
                    # No closing delimiter found, treat as regular text
                    current_text += delimiter + bold_text
                    
            # Check for italic (*text* or _text_)
            elif text[i] in ['*', '_'] and i > 0 and text[i-1] != text[i]:
                if current_text:
                    segments.append((current_text, False, False))
                    current_text = ""
                
                delimiter = text[i]
                i += 1
                italic_text = ""
                
                # Find closing delimiter
                while i < len(text):
                    if text[i] == delimiter:
                        segments.append((italic_text, False, True))
                        i += 1
                        break
                    italic_text += text[i]
                    i += 1
                else:
                    # No closing delimiter found, treat as regular text
                    current_text += delimiter + italic_text
            else:
                current_text += text[i]
                i += 1
        
        if current_text:
            segments.append((current_text, False, False))
        
        return segments
    
    def parse_markdown_file(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse markdown file and return structured content."""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        slides = []
        current_slide_content = {'title': '', 'bullets': [], 'text': '', 'needs_icon': False}
        
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            
            # Skip empty lines
            if not line:
                continue
                
            # Handle headers (new slides)
            if line.startswith('#'):
                # Save previous slide if it has content
                if current_slide_content['title'] or current_slide_content['bullets'] or current_slide_content['text']:
                    slides.append(current_slide_content.copy())
                
                # Start new slide
                title = re.sub(r'^#+\s*', '', line)
                # Determine if slide needs an icon based on title content
                needs_icon = any(keyword in title.lower() for keyword in 
                                ['machine learning', 'ai', 'data', 'business', 'strategy', 'roi', 'implementation'])
                current_slide_content = {'title': title, 'bullets': [], 'text': '', 'needs_icon': needs_icon}
                
            # Handle bullet points
            elif line.startswith('- ') or line.startswith('* '):
                bullet_text = re.sub(r'^[-*]\s*', '', line)
                current_slide_content['bullets'].append(bullet_text)
                
            # Handle regular text
            else:
                if current_slide_content['text']:
                    current_slide_content['text'] += '\n' + line
                else:
                    current_slide_content['text'] = line
        
        # Add the last slide
        if current_slide_content['title'] or current_slide_content['bullets'] or current_slide_content['text']:
            slides.append(current_slide_content)
            
        return slides
    
    def add_formatted_text(self, paragraph, text: str) -> None:
        """Add formatted text to a paragraph."""
        segments = self.parse_text_formatting(text)
        
        if not segments:
            return
            
        # Clear existing text
        paragraph.text = ""
        
        for i, (segment_text, is_bold, is_italic) in enumerate(segments):
            if i == 0:
                # Use the existing run
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.text = segment_text
            else:
                # Add new run
                run = paragraph.add_run()
                run.text = segment_text
            
            # Apply formatting
            run.font.bold = is_bold
            run.font.italic = is_italic
    
    def add_icon_shape(self, slide, slide_data: Dict[str, Any]) -> None:
        """Add a simple icon shape based on slide content."""
        if not slide_data.get('needs_icon', False):
            return
            
        # Determine icon type based on title
        title_lower = slide_data['title'].lower()
        
        # Add a simple geometric shape as icon
        left = Inches(8.5)
        top = Inches(1.5)
        width = Inches(1)
        height = Inches(1)
        
        if 'data' in title_lower or 'machine learning' in title_lower:
            # Add a circle for data/ML concepts
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(52, 152, 219)  # Blue
        elif 'business' in title_lower or 'roi' in title_lower:
            # Add a rectangle for business concepts
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(46, 204, 113)  # Green
        elif 'implementation' in title_lower or 'strategy' in title_lower:
            # Add a triangle for action/strategy concepts
            shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(230, 126, 34)  # Orange
        
        # Remove outline
        if 'shape' in locals():
            shape.line.color.rgb = RGBColor(255, 255, 255)
    
    def create_slide(self, slide_data: Dict[str, Any]) -> None:
        """Create a slide from parsed data."""
        # Choose layout based on content
        if slide_data['bullets']:
            # Use bullet layout
            slide_layout = self.prs.slide_layouts[1]  # Title and Content
        else:
            # Use title and content layout
            slide_layout = self.prs.slide_layouts[1]  # Title and Content
            
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title with formatting
        if slide_data['title']:
            title_shape = slide.shapes.title
            title_paragraph = title_shape.text_frame.paragraphs[0]
            self.add_formatted_text(title_paragraph, slide_data['title'])
            title_paragraph.font.size = Pt(32)
            title_paragraph.font.bold = True
        
        # Add icon if appropriate
        self.add_icon_shape(slide, slide_data)
        
        # Add content
        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            # Add bullets
            if slide_data['bullets']:
                for i, bullet in enumerate(slide_data['bullets']):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    self.add_formatted_text(p, bullet)
                    p.level = 0
                    p.font.size = Pt(18)
            
            # Add text content
            elif slide_data['text']:
                p = text_frame.paragraphs[0]
                self.add_formatted_text(p, slide_data['text'])
                p.font.size = Pt(18)
    
    def create_title_slide(self, title: str, subtitle: str = "") -> None:
        """Create a title slide."""
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        # Format title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Format subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    
    def convert(self, markdown_file: str, output_file: str, title: str = None, subtitle: str = "") -> None:
        """Convert markdown file to PPTX."""
        # Parse markdown
        slides_data = self.parse_markdown_file(markdown_file)
        
        if not slides_data:
            print("No content found in markdown file")
            return
        
        # Create title slide if title is provided
        if title:
            self.create_title_slide(title, subtitle)
        elif slides_data:
            # Use first slide title as presentation title
            first_slide = slides_data[0]
            if first_slide['title']:
                self.create_title_slide(first_slide['title'], subtitle)
                slides_data = slides_data[1:]  # Skip first slide as it's now the title
        
        # Create content slides
        for slide_data in slides_data:
            self.create_slide(slide_data)
        
        # Save presentation
        self.prs.save(output_file)
        print(f"Presentation saved as: {output_file}")


def main():
    """Main function to handle command line arguments."""
    parser = argparse.ArgumentParser(
        description="Convert Markdown files to PowerPoint presentations (POC)"
    )
    parser.add_argument(
        "input_file", 
        help="Input Markdown file path"
    )
    parser.add_argument(
        "-o", "--output", 
        default="output.pptx",
        help="Output PPTX file path (default: output.pptx)"
    )
    parser.add_argument(
        "-t", "--title",
        help="Presentation title for title slide"
    )
    parser.add_argument(
        "-s", "--subtitle",
        default="",
        help="Presentation subtitle for title slide"
    )
    
    args = parser.parse_args()
    
    # Validate input file
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: Input file '{args.input_file}' does not exist")
        return 1
    
    if not input_path.suffix.lower() in ['.md', '.markdown']:
        print(f"Warning: Input file doesn't have .md or .markdown extension")
    
    # Create converter and convert
    converter = MarkdownToPPTX()
    
    try:
        converter.convert(
            args.input_file, 
            args.output, 
            args.title, 
            args.subtitle
        )
        print("Conversion completed successfully!")
        return 0
    except Exception as e:
        print(f"Error during conversion: {e}")
        return 1


if __name__ == "__main__":
    exit(main())